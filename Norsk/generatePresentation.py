from pptx import Presentation

# Opprett en ny PowerPoint-presentasjon
prs = Presentation()

# Definer en funksjon for å legge til en slide
def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1: Tittel
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Novellanalyse: 'Heime-aleine-fest' og 'Kveld'"
slide1.placeholders[1].text = "Ditt navn\nDato"

# Slide 2: Introduksjon
add_slide("Introduksjon", "Kort presentasjon av novellene\nForfattere og kontekst\nFormålet med analysen")

# Slide 3: Sammendrag av 'Heime-aleine-fest'
add_slide("Sammendrag av 'Heime-aleine-fest'", "Steffen forbereder seg på nyttårsaften, reflekterer over familie og fest\nMiljø: Trygg barndomsheim, kald vinternatt")

# Slide 4: Sammendrag av 'Kveld'
add_slide("Sammendrag av 'Kveld'", "Sindre besøker sine besteforeldre og opplever varme og omsorg i en kald vinterkveld\nMiljø: Hjemmekoselig, preget av tradisjon")

# Slide 5: Personskildringer
add_slide("Sammenligning av Personskildringer", "Steffen: Ungdommelig refleksjon og uavhengighet\nSindre: Familie og tilhørighet\nForeldre vs. Besteforeldre: Kontraster i omsorg")

# Slide 6: Miljøskildringer
add_slide("Sammenligning av Miljøskildringer", "Trygghet vs. Selvstendighet\nSymbolikk i vinteren – kulde som kontrast til varme familiebånd")

# Slide 7: Likheter og Forskjeller
add_slide("Likheter og Forskjeller", "Begge noveller skildrer vinter og familiebånd\nSteffen opplever frigjøring, Sindre opplever tilknytning")

# Slide 8: Konklusjon
add_slide("Konklusjon", "Hvordan person- og miljøskildringer bygger opp novellenes budskap\nOvergangen mellom barndom, ungdom og familieforhold")

# Slide 9: Kilder
add_slide("Kilder", "Referanser til novellene og annen relevant litteratur")

# Lagre presentasjonen
prs.save("NovellanalysePrøvemuntlig.pptx")

print("PowerPoint-fil 'Novellanalyse.pptx' er generert!")
